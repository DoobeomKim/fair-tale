from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "fairtale-concept-deck.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(22, 31, 45),
    "muted": RGBColor(87, 98, 116),
    "line": RGBColor(223, 229, 238),
    "navy": RGBColor(21, 49, 86),
    "blue": RGBColor(44, 104, 196),
    "teal": RGBColor(18, 151, 139),
    "mint": RGBColor(230, 248, 244),
    "pale": RGBColor(247, 249, 252),
    "white": RGBColor(255, 255, 255),
    "gold": RGBColor(214, 163, 70),
    "soft_gold": RGBColor(255, 246, 224),
    "red": RGBColor(202, 80, 80),
}

FONT = "Pretendard"
FONT_FALLBACK = "Pretendard"


def set_text_style(run, size=20, bold=False, color="ink"):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(tag, namespaces=r_pr.nsmap)
        if el is None:
            from pptx.oxml import parse_xml

            el = parse_xml(f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{FONT}"/>')
            r_pr.append(el)
        else:
            el.set("typeface", FONT)


def add_textbox(slide, x, y, w, h, text="", size=22, bold=False, color="ink", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text_style(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.72, 0.48, 8.8, 0.45, "FAIRTALE", size=13, bold=True, color="teal")
    add_textbox(slide, 0.72, 0.93, 9.8, 0.72, title, size=30, bold=True, color="ink")
    if subtitle:
        add_textbox(slide, 0.75, 1.62, 10.7, 0.5, subtitle, size=17, color="muted")
    add_line(slide, 0.74, 2.18, 11.85, color="line")


def add_line(slide, x, y, w, color="line"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS[color]
    line.line.fill.background()
    return line


def add_round_rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    return shape


def add_card(slide, x, y, w, h, title, body, accent="blue"):
    add_round_rect(slide, x, y, w, h, fill="white", line="line")
    add_textbox(slide, x + 0.25, y + 0.22, w - 0.5, 0.35, title, size=17, bold=True, color=accent)
    tx = add_textbox(slide, x + 0.25, y + 0.68, w - 0.5, h - 0.85, body, size=13, color="ink")
    tx.text_frame.word_wrap = True
    return tx


def add_pill(slide, x, y, text, fill="mint", color="teal", width=2.5):
    shape = add_round_rect(slide, x, y, width, 0.42, fill=fill, line=fill)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_text_style(r, size=13, bold=True, color=color)
    return shape


def add_footer(slide, page):
    add_textbox(slide, 0.72, 7.08, 4.0, 0.22, "Fairtale Concept Deck", size=9, color="muted")
    add_textbox(slide, 12.0, 7.08, 0.6, 0.22, f"{page:02d}", size=9, color="muted", align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, size=17, color="ink", gap=0.13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap * 72)
        for run in p.runs:
            set_text_style(run, size=size, color=color)
    return box


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_cover(prs):
    s = blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["pale"]
    add_pill(s, 0.8, 0.72, "Cross-border B2B Expansion Platform", width=4.0)
    add_textbox(s, 0.78, 1.5, 8.8, 0.8, "Fairtale", size=44, bold=True, color="navy")
    add_textbox(s, 0.84, 2.28, 8.8, 1.15, "해외진출이 막막한 기업이\n바이어를 만나고 실제 기회로 나아가도록 돕는 플랫폼", size=24, bold=True, color="ink")
    add_textbox(s, 0.86, 3.72, 7.6, 0.7, "AI 통역은 기능입니다. Fairtale의 본질은 해외진출 성공 경험입니다.", size=16, color="muted")
    add_round_rect(s, 8.65, 1.0, 3.8, 4.95, fill="white", line="line")
    add_textbox(s, 9.05, 1.48, 3.0, 0.6, "North Star", size=18, bold=True, color="teal", align=PP_ALIGN.CENTER)
    add_textbox(s, 8.98, 2.28, 3.12, 1.9, "“Fairtale 덕분에\n해외 바이어를 만났고,\n실제 비즈니스 기회가\n생겼다.”", size=17, bold=True, color="ink", align=PP_ALIGN.CENTER)
    add_footer(s, 1)


def slide_problem(prs):
    s = blank(prs)
    add_title(s, "해외진출은 좋은 제품만으로 되지 않는다", "많은 중소기업은 제품은 있지만 시작점, 언어, 네트워크, 실무 절차에서 막힌다.")
    add_card(s, 0.78, 2.65, 2.85, 2.4, "시작점 부재", "어떤 시장에, 누구에게, 어떤 순서로 접근해야 할지 모른다.", "blue")
    add_card(s, 3.9, 2.65, 2.85, 2.4, "언어 장벽", "상담과 후속 커뮤니케이션에서 자신감이 떨어지고 속도가 느려진다.", "teal")
    add_card(s, 7.02, 2.65, 2.85, 2.4, "신뢰 네트워크 부족", "검증된 바이어, 수입사, 실무 파트너를 만나기 어렵다.", "gold")
    add_card(s, 10.14, 2.65, 2.45, 2.4, "실무 리스크", "계약, 인증, 통관, 물류, 현지화에서 다시 멈춘다.", "red")
    add_footer(s, 2)


def slide_essence(prs):
    s = blank(prs)
    add_title(s, "Fairtale은 AI 통역기가 아니다", "Fairtale은 해외진출을 실제 상담과 계약 기회까지 연결하는 실행 플랫폼이다.")
    add_round_rect(s, 0.9, 2.55, 5.35, 2.8, fill="soft_gold", line="soft_gold")
    add_textbox(s, 1.25, 2.95, 4.55, 0.5, "기술 중심 관점", size=19, bold=True, color="gold")
    add_bullets(s, 1.25, 3.58, 4.7, 1.1, ["실시간 음성통역", "화상회의", "자막 번역"], size=16)
    add_textbox(s, 6.55, 3.55, 0.7, 0.5, "→", size=30, bold=True, color="teal", align=PP_ALIGN.CENTER)
    add_round_rect(s, 7.45, 2.35, 4.85, 3.15, fill="mint", line="mint")
    add_textbox(s, 7.82, 2.85, 4.1, 0.55, "Fairtale 관점", size=20, bold=True, color="teal")
    add_bullets(s, 7.82, 3.52, 4.1, 1.45, ["해외진출 절차 안내", "적합한 바이어 연결", "상담 이후 실무 파트너 연결"], size=16)
    add_footer(s, 3)


def slide_journey(prs):
    s = blank(prs)
    add_title(s, "사용자 여정", "기업은 복잡한 해외진출 절차를 Fairtale이 안내하는 흐름대로 따라간다.")
    steps = [
        ("1", "등록", "회사/제품/인증/목표 시장"),
        ("2", "매칭", "적합한 바이어와 파트너 후보"),
        ("3", "상담", "화상 미팅과 실시간 번역 자막"),
        ("4", "정리", "상담 요약과 다음 액션"),
        ("5", "실행", "계약/인증/물류/마케팅 연결"),
    ]
    x = 0.8
    for num, title, body in steps:
        add_round_rect(s, x, 2.55, 2.25, 2.6, fill="white", line="line")
        add_pill(s, x + 0.25, 2.88, num, width=0.55)
        add_textbox(s, x + 0.25, 3.5, 1.7, 0.35, title, size=18, bold=True, color="ink")
        add_textbox(s, x + 0.25, 4.05, 1.7, 0.7, body, size=14, color="muted")
        x += 2.48
    add_footer(s, 4)


def slide_positioning(prs):
    s = blank(prs)
    add_title(s, "포지셔닝", "통역 기능이 있는 B2B 매칭 플랫폼이 아니라, 해외진출 성공 경험을 만드는 플랫폼.")
    add_round_rect(s, 1.0, 2.25, 11.3, 1.35, fill="mint", line="mint")
    add_textbox(s, 1.38, 2.6, 10.55, 0.8, "아시아 기업이 해외 바이어를 만나고, 언어 장벽 없이 상담하며,\n실제 수출 기회로 이어지도록 돕는다.", size=19, bold=True, color="ink", align=PP_ALIGN.CENTER)
    add_card(s, 1.0, 4.15, 3.35, 1.5, "공급 기업", "해외진출 절차를 쉽게 따라가고 바이어 상담 기회를 얻는다.", "blue")
    add_card(s, 4.98, 4.15, 3.35, 1.5, "해외 바이어", "검토할 가치가 있는 아시아 공급사를 더 쉽게 발견한다.", "teal")
    add_card(s, 8.95, 4.15, 3.35, 1.5, "파트너/투자자", "매칭, 온보딩, 언어 지원을 결합해 거래 마찰을 줄인다.", "gold")
    add_footer(s, 5)


def slide_moat(prs):
    s = blank(prs)
    add_title(s, "복제하기 어려운 자산", "기능은 따라 만들 수 있지만, 신뢰 네트워크와 거래 히스토리는 시간이 쌓여야 만들어진다.")
    add_card(s, 0.85, 2.45, 3.65, 2.3, "바이어 네트워크", "실제 구매 의향이 있는 해외 바이어, 수입사, 유통 파트너", "blue")
    add_card(s, 4.85, 2.45, 3.65, 2.3, "실무 파트너", "법무, 인증, 통관, 물류, 현지화, 마케팅 파트너", "teal")
    add_card(s, 8.85, 2.45, 3.65, 2.3, "거래 데이터", "어떤 제품이 어떤 시장에서 반응했는지에 대한 상담/거래 히스토리", "gold")
    add_round_rect(s, 1.35, 5.35, 10.6, 0.75, fill="pale", line="line")
    add_textbox(s, 1.65, 5.48, 10.0, 0.5, "Fairtale의 장기 방어력은 코드가 아니라,\n검증된 네트워크와 반복 가능한 해외진출 프로세스에서 나온다.", size=15, bold=True, color="ink", align=PP_ALIGN.CENTER)
    add_footer(s, 6)


def slide_partner_package(prs):
    s = blank(prs)
    add_title(s, "초기 파트너 패키지", "처음부터 큰 마켓플레이스보다, 해외진출 과정에서 자주 막히는 문제를 선별해서 해결한다.")
    items = [
        ("계약", "수출 계약서 1차 검토"),
        ("현지화", "제품 영문 소개서 및 카탈로그 현지화"),
        ("인증", "유럽/해외 인증 필요 여부 체크"),
        ("물류", "샘플 배송 및 물류 상담"),
        ("검증", "현지 바이어 또는 수입사 신뢰도 검토"),
    ]
    y = 2.28
    for label, body in items:
        add_pill(s, 1.08, y + 0.03, label, width=1.25)
        add_textbox(s, 2.62, y, 8.7, 0.38, body, size=18, bold=True, color="ink")
        add_line(s, 1.05, y + 0.58, 10.95, color="line")
        y += 0.78
    add_textbox(s, 1.08, 6.28, 10.7, 0.5, "운영 원칙: 적은 수의 검증된 파트너로 시작하고, 품질 관리와 책임 범위를 명확히 한다.", size=15, color="muted")
    add_footer(s, 7)


def slide_mvp(prs):
    s = blank(prs)
    add_title(s, "MVP 방향", "초기 MVP는 번역 기술보다 해외진출 여정이 작동한다는 것을 보여줘야 한다.")
    add_card(s, 0.9, 2.4, 3.5, 2.05, "1. 기업/제품 등록", "회사 정보, 제품, 인증, 희망 시장을 입력한다.", "blue")
    add_card(s, 4.85, 2.4, 3.5, 2.05, "2. 매칭/상담 요청", "카테고리나 시장 기반으로 바이어 상담 기회를 만든다.", "teal")
    add_card(s, 8.8, 2.4, 3.5, 2.05, "3. 미팅", "화상 미팅과 실시간 번역 자막으로 언어 장벽을 낮춘다.", "gold")
    add_round_rect(s, 2.3, 5.05, 8.75, 0.92, fill="pale", line="line")
    add_textbox(s, 2.65, 5.33, 8.05, 0.35, "TTS 음성통역은 안정성이 충분할 때 고도화 기능으로 확장한다.", size=16, bold=True, color="ink", align=PP_ALIGN.CENTER)
    add_footer(s, 8)


def slide_not(prs):
    s = blank(prs)
    add_title(s, "Fairtale이 되어서는 안 되는 것", "기능의 나열이 아니라, 해외진출 성공 경험을 만드는 서비스로 기억되어야 한다.")
    left = ["단순 화상회의 서비스", "단순 AI 번역기", "일반 온라인 전시 플랫폼", "기업 목록 디렉토리", "실시간 통역 기술 데모"]
    right = ["해외진출 안내자", "바이어 만남을 만드는 플랫폼", "상담 이후 실행까지 잇는 네트워크", "신뢰할 수 있는 파트너 연결", "실제 비즈니스 기회 창출"]
    add_round_rect(s, 0.9, 2.25, 5.1, 3.6, fill="soft_gold", line="soft_gold")
    add_textbox(s, 1.25, 2.62, 4.35, 0.4, "아닌 것", size=20, bold=True, color="red")
    add_bullets(s, 1.25, 3.25, 4.35, 1.9, left, size=16)
    add_round_rect(s, 7.05, 2.25, 5.1, 3.6, fill="mint", line="mint")
    add_textbox(s, 7.4, 2.62, 4.35, 0.4, "되어야 하는 것", size=20, bold=True, color="teal")
    add_bullets(s, 7.4, 3.25, 4.35, 1.9, right, size=16)
    add_footer(s, 9)


def slide_close(prs):
    s = blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]
    add_textbox(s, 0.95, 0.85, 3.5, 0.4, "FAIRTALE", size=14, bold=True, color="teal")
    add_textbox(s, 0.95, 1.85, 10.8, 1.2, "해외진출이 막막했던 기업에게\n첫 바이어와의 만남을 만들어준다.", size=32, bold=True, color="white")
    add_round_rect(s, 0.95, 4.15, 11.05, 1.2, fill="white", line="white")
    add_textbox(s, 1.35, 4.47, 10.2, 0.55, "“정말 감사합니다. 우리 회사에 정말 도움이 되었어요, Fairtale.”", size=20, bold=True, color="ink", align=PP_ALIGN.CENTER)
    add_textbox(s, 0.95, 6.58, 8.5, 0.35, "Concept deck based on docs/concept.md", size=10, color="white")


def build():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    for fn in [
        slide_cover,
        slide_problem,
        slide_essence,
        slide_journey,
        slide_positioning,
        slide_moat,
        slide_partner_package,
        slide_mvp,
        slide_not,
        slide_close,
    ]:
        fn(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
